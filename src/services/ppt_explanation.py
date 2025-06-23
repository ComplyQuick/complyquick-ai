# src/services/ppt_explanation.py
from pptx import Presentation
from .base_openai_service import BaseOpenAIService
from .storage_service import StorageService
from ..models import SlideExplanation
import logging
import os
import re
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np
from concurrent.futures import ThreadPoolExecutor, as_completed
import threading

logger = logging.getLogger(__name__)

class PPTExplanationService(BaseOpenAIService):
    def __init__(self):
        super().__init__()
        self.storage_service = StorageService()
        # Initialize sentence transformer for semantic similarity
        try:
            self.sentence_model = SentenceTransformer('all-MiniLM-L6-v2')
            logger.info("Sentence transformer model loaded successfully")
        except Exception as e:
            logger.warning(f"Could not load sentence transformer model: {e}")
            self.sentence_model = None

    def extract_slide_content(self, ppt_path: str):
        """
        Extract text content from each slide in the PowerPoint presentation.
        """
        logger.info(f"Extracting content from PPT: {ppt_path}")
        presentation = Presentation(ppt_path)
        slides_content = []

        for slide_number, slide in enumerate(presentation.slides, 1):
            slide_text = []
            
            # Extract text from all shapes including text boxes, titles, and content
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Clean and preserve the text structure
                    text = shape.text.strip()
                    # Preserve bullet points and numbering
                    if text.startswith(('•', '-', '*', '1.', '2.', '3.', '4.', '5.', '6.', '7.', '8.', '9.')):
                        slide_text.append(text)
                    else:
                        slide_text.append(text)
                
                # Handle table content if present
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_text.append(" | ".join(row_text))
            
            # Join all content while preserving structure
            content = " ".join(slide_text)
            slides_content.append(content)
            
            # Log detailed content for debugging
            logger.debug(f"Extracted content from slide {slide_number}: {content[:100]}...")

        logger.info(f"Extracted content from {len(slides_content)} slides")
        return slides_content

    

    def _create_prompt(self, index: int, total_slides: int, slide_text: str, company_name: str, pocs: list, relevant_contacts: dict) -> str:
        """
        Create the appropriate prompt based on slide position and content.
        """
        # Create role introductions and personalized examples
        role_introductions = []
        personalized_examples = []
        company_examples = []
        
        for poc in pocs:
            role_lower = poc['role'].lower()
            # First introduce the person with their role
            role_introductions.append(f"{poc['name']} as {poc['role']}")
            
            # Create innovative examples based on roles
            if 'ceo' in role_lower:
                personalized_examples.append(f"{poc['name']} champions our culture of innovation and excellence")
                company_examples.append(f"Under {poc['name']}'s leadership, {company_name} has set new industry standards")
            elif 'cto' in role_lower:
                personalized_examples.append(f"{poc['name']} spearheads our digital transformation initiatives")
                company_examples.append(f"{company_name}'s technological edge is driven by {poc['name']}'s vision")
            elif 'hr' in role_lower:
                personalized_examples.append(f"{poc['name']} cultivates our diverse and inclusive workplace culture")
                company_examples.append(f"Thanks to {poc['name']}'s initiatives, {company_name} has become a workplace of choice")
            elif 'risk' in role_lower:
                personalized_examples.append(f"{poc['name']} ensures our robust security framework")
                company_examples.append(f"{company_name}'s security protocols, led by {poc['name']}, are industry-leading")
            elif 'legal' in role_lower:
                personalized_examples.append(f"{poc['name']} upholds our commitment to compliance and ethics")
                company_examples.append(f"{company_name}'s compliance standards, guided by {poc['name']}, set the benchmark")
            else:
                personalized_examples.append(f"{poc['name']} drives excellence in their domain")
                company_examples.append(f"{company_name}'s success is amplified by {poc['name']}'s contributions")

        # Create a dynamic leadership section with varied language
        leadership_section = f"At {company_name}, we're proud to have " + ", ".join(role_introductions) + ". " + ", ".join(personalized_examples) + ". " + ", ".join(company_examples) + "."

        # CRITICAL: Add comprehensive coverage instruction for all slides
        comprehensive_instruction = (
            "IMPORTANT: You MUST cover EVERY single point, detail, and piece of information mentioned in the slide content. "
            "Do not skip any bullet points, numbers, percentages, definitions, or any other content. "
            "Your explanation should be comprehensive and thorough, ensuring that every element from the slide is addressed. "
            "If the slide contains lists, cover each item. If it has statistics, mention them. If it has definitions, explain them. "
            "The goal is to create a complete explanation that leaves nothing out while maintaining a friendly, conversational tone. "
            "Before writing your explanation, mentally check off each point from the slide to ensure nothing is missed. "
            "If you see bullet points, numbers, percentages, or any structured content, make sure to address each one specifically. "
            "Your explanation should be detailed enough that someone who hasn't seen the slide would understand every point mentioned. "
            "STRUCTURE FOR TEXT-TO-SPEECH: Write in a natural, conversational style that flows well when spoken aloud. "
            "Use smooth transitions between ideas, avoid abrupt topic changes, and create a natural rhythm. "
            "Break complex information into digestible chunks with natural pauses. "
            "Use connecting words like 'furthermore', 'additionally', 'moreover', 'in addition', 'also', 'besides', 'similarly', 'likewise' to create flow. "
            "Avoid long, run-on sentences that are difficult to follow when spoken. "
            "Use varied sentence structures and natural speech patterns that sound human and engaging. "
            "FORMATTING: Do NOT use any markdown formatting, double asterisks (**), bold text, or any special characters that could interfere with text-to-speech systems. "
            "Write in plain text only, suitable for natural speech synthesis."
        )

        if index == 0:
            # First slide
            prompt = (
                f"Create a warm, engaging introduction for {company_name}'s training presentation. "
                f"The explanation should be natural and suitable for text-to-speech narration. "
                f"Make it friendly and professional, about 4-5 sentences long. "
                f"Don't use time-specific greetings like 'good morning/evening'. "
                f"Focus on welcoming the audience and introducing the topic. "
                f"Write in a conversational tone that flows naturally when spoken aloud. "
                f"Use smooth transitions and natural speech patterns. "
                f"{leadership_section} "
                f"Emphasize how {company_name} is leading the way in this field. "
                f"Incorporate specific examples of how our executives have championed this initiative. "
                f"{comprehensive_instruction}\n\n"
                f"Slide content:\n{slide_text}"
            )
        elif index == total_slides - 1:
            # Last slide
            poc_names = ", ".join([poc['name'] for poc in pocs])
            prompt = (
                f"Create a friendly conclusion for {company_name}'s training presentation. "
                f"The explanation should be natural for text-to-speech narration. "
                f"Keep it warm and professional, about 3-4 sentences. "
                f"Write in a conversational tone that flows naturally when spoken aloud. "
                f"Use smooth transitions and natural speech patterns. "
                f"Highlight how {poc_names} and their teams "
                f"embody these principles in their daily work. "
                f"Share a specific example of how {company_name} is implementing these practices. "
                f"{leadership_section} "
                f"Don't use phrases like 'thank you for your time today'. "
                f"{comprehensive_instruction}\n\n"
                f"Slide content:\n{slide_text}"
            )
        else:
            # Regular slides
            # Determine the appropriate tone and examples based on content
            tone_instruction = ""
            role_examples = []
            
            if any(keyword in slide_text.lower() for keyword in ['harassment', 'posh', 'sexual harassment']):
                tone_instruction = (
                    "Use a professional and supportive tone. "
                    "Emphasize how our HR and Legal departments, "
                    "actively promote a safe and respectful workplace. "
                )
                for poc in pocs:
                    if 'hr' in poc['role'].lower() or 'legal' in poc['role'].lower():
                        role_examples.append(f"{poc['name']} regularly conducts awareness sessions")
            elif any(keyword in slide_text.lower() for keyword in ['security', 'compliance', 'data protection']):
                tone_instruction = (
                    "Use a professional and authoritative tone. "
                    "Reference how our technology and risk management teams "
                    "maintain our security standards. "
                )
                for poc in pocs:
                    if 'cto' in poc['role'].lower() or 'risk' in poc['role'].lower():
                        role_examples.append(f"{poc['name']} implements cutting-edge security measures")
            elif any(keyword in slide_text.lower() for keyword in ['benefits', 'leave', 'policy', 'employee']):
                tone_instruction = (
                    "Use a warm and informative tone. "
                    "Highlight how our HR department "
                    "ensures employee well-being. "
                )
                for poc in pocs:
                    if 'hr' in poc['role'].lower():
                        role_examples.append(f"{poc['name']} personally reviews our employee benefits program")
            elif any(keyword in slide_text.lower() for keyword in ['training', 'development', 'learning']):
                tone_instruction = (
                    "Use an encouraging and motivational tone. "
                    "Showcase how our executives and department heads "
                    "invest in employee growth. "
                )
                for poc in pocs:
                    if 'ceo' in poc['role'].lower():
                        role_examples.append(f"{poc['name']} champions our learning culture")
            
            role_examples_text = ", ".join(role_examples) if role_examples else "our company's commitment"
            prompt = (
                f"Create an engaging explanation of this slide for {company_name}'s training. "
                f"The explanation should be natural and suitable for text-to-speech narration. "
                f"Make it comprehensive and detailed, ensuring you cover ALL content from the slide. "
                f"Use a conversational tone while maintaining professionalism. "
                f"Break down complex concepts into clear explanations. "
                f"Write in a natural, flowing style that sounds great when spoken aloud. "
                f"Use varied sentence lengths and natural speech patterns. "
                f"Create smooth transitions between ideas and concepts. "
                f"Break information into digestible chunks with natural pauses. "
                f"{tone_instruction}"
                f"{leadership_section} "
                f"Include specific examples of how {company_name} implements these practices, "
                f"such as {role_examples_text}. "
                f"Avoid using bullet points or lists - structure the content in flowing paragraphs. "
                f"Don't use transition phrases like 'moving on' or 'in this slide'. "
                f"Instead, use natural connecting words like 'furthermore', 'additionally', 'moreover', 'in addition', 'also', 'besides', 'similarly', 'likewise'. "
                f"{comprehensive_instruction} "
                f"Structure your response to systematically cover each point mentioned in the slide content, "
                f"ensuring nothing is overlooked while maintaining a friendly, approachable tone that flows naturally when read aloud.\n\n"
                f"Slide content:\n{slide_text}"
            )
        
        return prompt

    def _verify_content_coverage(self, original_content: str, explanation: str) -> bool:
        """
        Verify that the explanation covers all key points from the original content.
        """
        try:
            # Extract key terms, numbers, and bullet points from original content
            # Find numbers, percentages, bullet points - use simpler patterns
            numbers = re.findall(r'\d+', original_content)  # Simplified number pattern
            bullets = re.findall(r'[•\-*]\s*([^•\-*\n]+)', original_content)
            
            logger.debug(f"Found {len(numbers)} numbers and {len(bullets)} bullet points in content")
            
            # Check if numbers are mentioned in explanation
            numbers_covered = True
            if numbers:
                # Check if each number is mentioned in the explanation
                missing_numbers = []
                for num in numbers:
                    if num not in explanation:
                        missing_numbers.append(num)
                numbers_covered = len(missing_numbers) == 0
                
                # Log missing numbers for debugging
                if missing_numbers:
                    logger.warning(f"Missing numbers in explanation: {missing_numbers}")
            
            # Check if bullet points are covered (simplified check)
            bullets_covered = True
            if bullets:
                # Extract key terms from bullets
                bullet_terms = []
                for bullet in bullets:
                    # Extract first few important words from each bullet
                    words = bullet.strip().split()[:3]
                    bullet_terms.extend(words)
                
                # Check if most bullet terms are mentioned
                covered_terms = 0
                for term in bullet_terms:
                    if term.lower() in explanation.lower():
                        covered_terms += 1
                
                bullets_covered = covered_terms >= len(bullet_terms) * 0.7  # 70% coverage threshold
                
                # Log coverage statistics
                coverage_percentage = (covered_terms / len(bullet_terms)) * 100 if bullet_terms else 100
                logger.debug(f"Bullet point coverage: {coverage_percentage:.1f}% ({covered_terms}/{len(bullet_terms)})")
                
                if not bullets_covered:
                    missing_terms = [term for term in bullet_terms if term.lower() not in explanation.lower()]
                    logger.warning(f"Missing bullet point terms: {missing_terms[:5]}...")  # Show first 5 missing terms
            
            result = numbers_covered and bullets_covered
            logger.debug(f"Content coverage verification result: {result}")
            return result
            
        except Exception as e:
            logger.error(f"Error in _verify_content_coverage: {str(e)}", exc_info=True)
            # Return True to avoid blocking the process if verification fails
            return True

    def _semantic_verify_content_coverage(self, original_content: str, explanation: str) -> bool:
        """
        Verify content coverage using semantic similarity with sentence transformers.
        """
        if self.sentence_model is None:
            logger.warning("Sentence transformer not available, falling back to regex verification")
            return self._verify_content_coverage(original_content, explanation)
        
        try:
            # Extract key content segments from original slide
            content_segments = self._extract_content_segments(original_content)
            
            if not content_segments:
                logger.debug("No content segments found, skipping semantic verification")
                return True
            
            # Calculate semantic similarity for each segment
            similarities = []
            for segment in content_segments:
                if segment.strip():
                    # Get embeddings for segment and explanation
                    segment_embedding = self.sentence_model.encode([segment])
                    explanation_embedding = self.sentence_model.encode([explanation])
                    
                    # Calculate cosine similarity
                    similarity = cosine_similarity(segment_embedding, explanation_embedding)[0][0]
                    similarities.append(similarity)
                    
                    logger.debug(f"Segment: '{segment[:50]}...' - Similarity: {similarity:.3f}")
            
            # Calculate average similarity and coverage
            avg_similarity = np.mean(similarities) if similarities else 0
            coverage_threshold = 0.6  # 60% similarity threshold
            
            logger.debug(f"Average semantic similarity: {avg_similarity:.3f}")
            logger.debug(f"Coverage threshold: {coverage_threshold}")
            
            return avg_similarity >= coverage_threshold
            
        except Exception as e:
            logger.error(f"Error in semantic verification: {str(e)}", exc_info=True)
            # Fall back to regex verification
            return self._verify_content_coverage(original_content, explanation)
    
    def _extract_content_segments(self, content: str) -> list:
        """
        Extract meaningful content segments from slide content for semantic verification.
        """
        segments = []
        
        # Split by bullet points, numbers, and other delimiters
        lines = content.split('\n')
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Remove bullet points and numbering
            cleaned_line = re.sub(r'^[•\-*\d\.\s]+', '', line)
            if cleaned_line and len(cleaned_line) > 3:  # Minimum meaningful length
                segments.append(cleaned_line)
        
        # If no segments found, split by sentences
        if not segments:
            sentences = re.split(r'[.!?]+', content)
            segments = [s.strip() for s in sentences if s.strip() and len(s.strip()) > 10]
        
        # If still no segments, use the entire content
        if not segments:
            segments = [content]
        
        logger.debug(f"Extracted {len(segments)} content segments for semantic verification")
        return segments

    def _process_single_slide(self, args):
        """
        Process a single slide with explanation generation and verification.
        This method is designed to be used with concurrent processing.
        """
        index, slide_text, total_slides, company_name, pocs = args
        
        try:
            logger.info(f"Processing slide {index + 1}/{total_slides}")
            
            # Step 1: Create prompt
            logger.debug(f"Creating prompt for slide {index + 1}")
            prompt = self._create_prompt(index, total_slides, slide_text, company_name, pocs, {})
            
            # Step 2: Make OpenAI request
            logger.debug(f"Making OpenAI request for slide {index + 1}")
            explanation = self._make_openai_request(prompt)
            
            # Step 3: Clean explanation
            logger.debug(f"Cleaning explanation for slide {index + 1}")
            try:
                if explanation is None:
                    explanation = "No explanation generated"
                elif not isinstance(explanation, str):
                    explanation = str(explanation)
                
                cleaned_explanation = (explanation
                    .replace("•", "")
                    .replace("-", "")
                    .replace("...", ".")
                    .replace("\n\n", " ")
                    .replace("**", "")  # Remove double asterisks to prevent double highlighting
                    .strip())
            except Exception as cleaning_error:
                logger.error(f"Error cleaning explanation for slide {index + 1}: {str(cleaning_error)}")
                cleaned_explanation = str(explanation) if explanation else "Error processing explanation"
            
            # Step 4: Verify content coverage and regenerate if needed
            logger.debug(f"Verifying content coverage for slide {index + 1}")
            verification_enabled = True
            
            if verification_enabled:
                try:
                    if not self._semantic_verify_content_coverage(slide_text, cleaned_explanation):
                        logger.warning(f"Content coverage insufficient for slide {index + 1}, regenerating...")
                        enhanced_prompt = (
                            f"CRITICAL: The previous explanation missed some content. "
                            f"Please create a COMPLETE explanation that covers EVERY single point from this slide content. "
                            f"Make sure to mention all numbers, percentages, bullet points, and any other details. "
                            f"Be thorough and comprehensive while maintaining a friendly tone. "
                            f"Write in a natural, conversational style that flows well when spoken aloud. "
                            f"Use smooth transitions between ideas and create a natural rhythm. "
                            f"Break complex information into digestible chunks with natural pauses. "
                            f"Use connecting words like 'furthermore', 'additionally', 'moreover', 'in addition', 'also', 'besides', 'similarly', 'likewise' to create flow. "
                            f"Avoid long, run-on sentences that are difficult to follow when spoken. "
                            f"Use varied sentence structures and natural speech patterns. "
                            f"Original slide content:\n{slide_text}\n\n"
                            f"Previous incomplete explanation:\n{cleaned_explanation}\n\n"
                            f"Please provide a complete explanation that covers everything in a natural, TTS-friendly style:"
                        )
                        
                        enhanced_explanation = self._make_openai_request(enhanced_prompt)
                        if enhanced_explanation is None:
                            enhanced_explanation = "No enhanced explanation generated"
                        elif not isinstance(enhanced_explanation, str):
                            enhanced_explanation = str(enhanced_explanation)
                        
                        cleaned_explanation = (enhanced_explanation
                            .replace("•", "")
                            .replace("-", "")
                            .replace("...", ".")
                            .replace("\n\n", " ")
                            .replace("**", "")
                            .strip())
                except Exception as verification_error:
                    logger.warning(f"Content verification failed for slide {index + 1}, continuing with original explanation: {str(verification_error)}")
            
            logger.debug(f"Generated explanation for slide {index + 1}: {cleaned_explanation[:200]}...")
            return index, cleaned_explanation
            
        except Exception as e:
            logger.error(f"Error processing slide {index + 1}: {str(e)}", exc_info=True)
            return index, f"Error generating explanation: {str(e)}"

    def generate_explanations(self, slides_content: list, company_name: str, pocs: list):
        """
        Generate explanations for each slide using the OpenAI API with concurrent processing.
        """
        logger.info(f"Generating explanations for {len(slides_content)} slides with concurrent processing")
        logger.info(f"Company name: {company_name}")
        
        total_slides = len(slides_content)
        
        # Determine optimal number of workers based on content length
        # More workers for shorter content, fewer for longer content to avoid rate limits
        avg_content_length = sum(len(content) for content in slides_content) / len(slides_content)
        if avg_content_length > 500:
            max_workers = 3  # Fewer workers for longer content to avoid rate limits
        elif avg_content_length > 200:
            max_workers = 5  # Medium workers for medium content
        else:
            max_workers = 8  # More workers for shorter content
        
        logger.info(f"Using {max_workers} concurrent workers for processing")
        
        # Prepare arguments for concurrent processing
        args_list = [(index, slide_text, total_slides, company_name, pocs) 
                    for index, slide_text in enumerate(slides_content)]
        
        explanations = [None] * total_slides  # Pre-allocate list
        
        # Process slides concurrently
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            # Submit all tasks
            future_to_index = {executor.submit(self._process_single_slide, args): args[0] 
                             for args in args_list}
            
            # Collect results as they complete
            completed_count = 0
            for future in as_completed(future_to_index):
                try:
                    index, explanation = future.result()
                    explanations[index] = explanation
                    completed_count += 1
                    logger.info(f"Completed {completed_count}/{total_slides} slides")
                except Exception as e:
                    index = future_to_index[future]
                    logger.error(f"Error in concurrent processing for slide {index + 1}: {str(e)}")
                    explanations[index] = f"Error generating explanation: {str(e)}"
                    completed_count += 1

        logger.info(f"Completed generating {len(explanations)} explanations with concurrent processing")
        return explanations

    def process_ppt(self, presentation_url: str, company_name: str, pocs: list) -> list:
        """
        Process the PPT to extract content and generate explanations for each slide.
        """
        logger.info(f"Starting PPT processing for {company_name}")
        try:
            # Use storage_service to download and get the file path
            ppt_path = self.storage_service.download_presentation(presentation_url)
            slides_content = self.extract_slide_content(ppt_path)
            explanations = self.generate_explanations(slides_content, company_name, pocs)
            
            result = []
            for i in range(len(slides_content)):
                # Generate a concise gist for the content
                gist_prompt = (
                    f"Create a concise title or gist for the following slide content. "
                    f"Focus on the main topic or key message. Be brief and direct, but ensure it captures the complete topic:\n\n"
                    f"{slides_content[i]}"
                )
                gist = self._make_openai_request(gist_prompt)
                cleaned_gist = gist.strip().replace('"', '').replace("'", '')
                
                # If the gist is too long, try to make it more concise
                if len(cleaned_gist.split()) > 8:
                    gist_prompt = (
                        f"Create a very concise title (4-6 words) for the following slide content. "
                        f"Focus on the main topic only:\n\n"
                    f"{slides_content[i]}"
                )
                gist = self._make_openai_request(gist_prompt)
                cleaned_gist = gist.strip().replace('"', '').replace("'", '')
                
                result.append(
                    SlideExplanation(
                        slide=i + 1,
                        content=cleaned_gist,
                        explanation=explanations[i]
                    )
                )
            
            # Clean up the downloaded file
            try:
                os.remove(ppt_path)
            except:
                pass  
            
            logger.info(f"Successfully processed PPT with {len(result)} slides")
            return result
        except Exception as e:
            logger.error("Error in process_ppt", exc_info=True)
            raise