import boto3
import os
from urllib.parse import urlparse, unquote_plus
from pptx import Presentation
import logging
from google.oauth2.credentials import Credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
import io
from dotenv import load_dotenv

logger = logging.getLogger(__name__)

class StorageService:
    def __init__(self):
        self.SCOPES = [
            'https://www.googleapis.com/auth/drive.readonly',
            'https://www.googleapis.com/auth/drive.file'
        ]
        self._initialize_google_credentials()
        
        # Validate Linode Object Storage credentials
        linode_access_key = os.getenv("LINODE_ACCESS_KEY")
        linode_secret_key = os.getenv("LINODE_SECRET_KEY")
        linode_region = os.getenv("LINODE_REGION")
        linode_endpoint = os.getenv("LINODE_ENDPOINT")
        self.bucket_name = os.getenv("LINODE_BUCKET_NAME")
        self.audio_folder = os.getenv("LINODE_AUDIO_FOLDER", "test-uploads")  # Default to test-uploads
        
        if not all([linode_access_key, linode_secret_key, linode_region, linode_endpoint, self.bucket_name]):
            logger.error("Missing Linode Object Storage credentials. Please check your environment variables.")
            raise ValueError("Linode credentials not properly configured. Required: LINODE_ACCESS_KEY, LINODE_SECRET_KEY, LINODE_REGION, LINODE_ENDPOINT, LINODE_BUCKET_NAME")
        
        try:
            self.s3_client = boto3.client(
                's3',
                aws_access_key_id=linode_access_key,
                aws_secret_access_key=linode_secret_key,
                region_name=linode_region,
                endpoint_url=linode_endpoint
            )
            # Test the credentials by making a simple S3 call
            self.s3_client.list_buckets()
            logger.info("Successfully initialized Linode Object Storage client")
        except Exception as e:
            logger.error(f"Failed to initialize Linode Object Storage client: {str(e)}")
            raise ValueError(f"Failed to initialize Linode Object Storage client: {str(e)}")

    def _initialize_google_credentials(self):
        """Initialize Google Drive credentials from environment variables"""
        try:
            self.creds = Credentials(
                None,  # Token is not needed as we'll use refresh token
                client_id=os.getenv('GOOGLE_DRIVE_CLIENT_ID'),
                client_secret=os.getenv('GOOGLE_DRIVE_CLIENT_SECRET'),
                token_uri='https://oauth2.googleapis.com/token',
                refresh_token=os.getenv('GOOGLE_DRIVE_REFRESH_TOKEN')
            )
            logger.info("Successfully initialized Google Drive credentials")
        except Exception as e:
            logger.error(f"Error initializing Google Drive credentials: {str(e)}")
            raise ValueError("Failed to initialize Google Drive credentials. Check your environment variables.")

    def _get_file_id_from_url(self, url: str) -> str:
        """Extract file ID from Google Docs URL"""
        try:
            parsed_url = urlparse(url)
            
            # Handle different Google Drive URL formats
            if 'docs.google.com' in parsed_url.netloc:
                path_parts = parsed_url.path.split('/')
                # The ID is usually after /d/ in the URL
                for i, part in enumerate(path_parts):
                    if part == 'd':
                        return path_parts[i + 1].split('/')[0]
            
            # Handle direct drive.google.com URLs
            if 'drive.google.com' in parsed_url.netloc:
                if 'id=' in url:
                    # Handle old-style URLs with id parameter
                    return parsed_url.query.split('id=')[1].split('&')[0]
                else:
                    # Handle new-style URLs
                    path_parts = parsed_url.path.split('/')
                    for i, part in enumerate(path_parts):
                        if part in ['d', 'file']:
                            return path_parts[i + 1]
            
            raise ValueError("Could not extract file ID from URL")
        except Exception as e:
            logger.error(f"Error extracting file ID from URL: {str(e)}")
            raise ValueError(f"Invalid Google Drive URL format: {url}")

    def download_presentation(self, presentation_url: str, download_path: str = None) -> str:
        """
        Download a presentation from either Google Drive or S3
        """
        try:
            logger.info(f"Starting download of presentation from URL: {presentation_url}")
            
            # If no download path is provided, extract extension from URL
            if not download_path:
                file_extension = os.path.splitext(presentation_url)[1]
                if not file_extension:
                    file_extension = '.pptx'  # Default to .pptx if no extension found
                download_path = f"downloaded_presentation{file_extension}"
            
            logger.info(f"Using download path: {download_path}")
            
            # Check if it's a Google Drive URL
            if 'docs.google.com' in presentation_url or 'drive.google.com' in presentation_url:
                logger.info("Detected Google Drive URL")
                return self._download_from_google_drive(presentation_url, download_path)
            # Check if it's an S3 compatible URL (AWS S3 or Linode Object Storage)
            elif 's3.' in presentation_url or '.amazonaws.com' in presentation_url or 'linodeobjects.com' in presentation_url:
                logger.info("Detected S3-compatible URL (AWS S3 or Linode Object Storage)")
                return self.download_ppt_from_s3(presentation_url, download_path)
            else:
                raise ValueError(f"Unsupported URL format: {presentation_url}. Must be either Google Drive or S3-compatible URL.")
                
        except Exception as e:
            logger.error(f"Error downloading presentation: {str(e)}", exc_info=True)
            raise Exception(f"Failed to download presentation: {str(e)}")

    def _download_from_google_drive(self, url: str, download_path: str) -> str:
        """Download presentation from Google Drive"""
        try:
            file_id = self._get_file_id_from_url(url)
            if not file_id:
                raise ValueError("Invalid Google Drive URL")

            service = build('drive', 'v3', credentials=self.creds)
            
            # Get the file metadata
            file_metadata = service.files().get(fileId=file_id, fields='mimeType').execute()
            mime_type = file_metadata.get('mimeType', '')
            
            # Different handling based on file type
            if mime_type == 'application/vnd.google-apps.presentation':
                # Native Google Slides
                request = service.files().export_media(
                    fileId=file_id,
                    mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
                )
            else:
                # Uploaded PowerPoint file
                request = service.files().get_media(fileId=file_id)
            
            file_handle = io.BytesIO()
            downloader = MediaIoBaseDownload(file_handle, request)
            done = False
            
            while not done:
                status, done = downloader.next_chunk()
                if status:
                    logger.info(f"Download progress: {int(status.progress() * 100)}%")

            file_handle.seek(0)
            with open(download_path, 'wb') as f:
                f.write(file_handle.read())

            logger.info(f"Successfully downloaded presentation to: {download_path}")
            return download_path

        except Exception as e:
            logger.error(f"Error downloading from Google Drive: {str(e)}", exc_info=True)
            raise Exception(f"Failed to download from Google Drive: {str(e)}")

    def download_ppt_from_s3(self, s3_url: str, download_path: str = "downloaded_presentation.pptx"):
        """
        Download a PowerPoint file from the given S3-compatible URL (AWS S3 or Linode Object Storage) and save it locally.
        """
        try:
            logger.info(f"Attempting to download from S3-compatible storage: {s3_url}")
            
            # Parse the S3-compatible URL
            parsed_url = urlparse(s3_url)
            
            # Try to extract bucket name from URL, but default to environment variable
            if 'linodeobjects.com' in parsed_url.netloc:
                # For Linode, use the bucket name from environment
                bucket_name = self.bucket_name
            else:
                # For AWS S3, extract from URL
                bucket_name = parsed_url.netloc.split('.')[0]
            
            # Fix double encoding issue
            key = parsed_url.path.lstrip('/')
            # First, decode %25 to % if it exists
            key = key.replace('%25', '%')
            # Then decode the remaining URL encoding
            key = unquote_plus(key)
            
            logger.info(f"Parsed S3-compatible URL - Bucket: {bucket_name}, Key: {key}")
            
            # Check if bucket exists
            try:
                self.s3_client.head_bucket(Bucket=bucket_name)
            except self.s3_client.exceptions.ClientError as e:
                error_code = e.response['Error']['Code']
                if error_code == '404':
                    raise Exception(f"Bucket '{bucket_name}' does not exist")
                elif error_code == '403':
                    raise Exception(f"Access denied to bucket '{bucket_name}'. Please check your storage credentials and permissions.")
                else:
                    raise Exception(f"Error accessing bucket '{bucket_name}': {str(e)}")

            # Check if object exists and is accessible
            try:
                self.s3_client.head_object(Bucket=bucket_name, Key=key)
            except self.s3_client.exceptions.ClientError as e:
                error_code = e.response['Error']['Code']
                if error_code == '404':
                    raise Exception(f"File '{key}' not found in bucket '{bucket_name}'")
                elif error_code == '403':
                    raise Exception(f"Access denied to file '{key}' in bucket '{bucket_name}'. Please check your storage credentials and permissions.")
                else:
                    raise Exception(f"Error accessing file '{key}' in bucket '{bucket_name}': {str(e)}")

            logger.info(f"Attempting to download file...")
            # Download the file
            self.s3_client.download_file(bucket_name, key, download_path)
            logger.info(f"Successfully downloaded file to: {download_path}")
            
            return download_path
            
        except self.s3_client.exceptions.NoSuchKey:
            logger.error(f"File not found in storage: {key}", exc_info=True)
            raise Exception(f"The object '{key}' does not exist in bucket '{bucket_name}'.")
        except Exception as e:
            logger.error(f"Error downloading from storage: {str(e)}", exc_info=True)
            raise Exception(f"Failed to download file from storage: {str(e)}")

    def extract_content_from_ppt(self, presentation_url: str):
        """
        Extract text content from a PowerPoint file from either Google Drive or S3-compatible storage
        """
        try:
            logger.info(f"Starting content extraction from presentation URL: {presentation_url}")
            
            # Use the generic download method instead of directly calling S3
            logger.info("Downloading presentation...")
            ppt_path = self.download_presentation(presentation_url)
            logger.info(f"Presentation downloaded to: {ppt_path}")
            
            logger.info("Opening presentation file...")
            presentation = Presentation(ppt_path)
            logger.info(f"Successfully opened presentation with {len(presentation.slides)} slides")
            
            structured_content = []
            
            for slide_number, slide in enumerate(presentation.slides, 1):
                logger.info(f"Processing slide {slide_number}")
                slide_content = {
                    "slide_number": slide_number,
                    "title": "",
                    "content": [],
                    "notes": ""
                }
                
                # Extract slide title
                if slide.shapes.title:
                    slide_content["title"] = slide.shapes.title.text.strip()
                    logger.info(f"Extracted title: {slide_content['title']}")
                
                # Extract text from shapes (bullet points, text boxes, etc.)
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if shape != slide.shapes.title:  # Skip title as we already got it
                            content = shape.text.strip()
                            slide_content["content"].append(content)
                            logger.info(f"Extracted content: {content[:100]}...")
                
                # Extract speaker notes
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    slide_content["notes"] = slide.notes_slide.notes_text_frame.text.strip()
                    logger.info(f"Extracted notes: {slide_content['notes'][:100]}...")
                
                structured_content.append(slide_content)
            
            # Create a comprehensive knowledge base
            logger.info("Creating knowledge base from structured content...")
            knowledge_base = self._create_knowledge_base(structured_content)
            logger.info(f"Knowledge base created with length: {len(knowledge_base)}")
            
            # Clean up the downloaded file
            try:
                logger.info("Cleaning up downloaded file...")
                os.remove(ppt_path)
                logger.info("File cleanup successful")
            except Exception as e:
                logger.warning(f"Failed to clean up file: {str(e)}")
            
            return knowledge_base
        except Exception as e:
            logger.error(f"Error extracting content from presentation: {str(e)}", exc_info=True)
            raise Exception(f"Failed to extract content from presentation: {str(e)}")

    def _create_knowledge_base(self, structured_content):
        """
        Convert structured PPT content into a detailed knowledge base format.
        """
        knowledge_base = []
        
        for slide in structured_content:
            # Start with the topic/title
            section = f"Topic: {slide['title']}\n\n"
            
            # Add main content points with context
            section += "Key Points:\n"
            for point in slide['content']:
                # Clean up bullet points and formatting
                clean_point = point.replace('•', '').strip()
                if clean_point:
                    section += f"- {clean_point}\n"
            
            # Add detailed explanation from notes if available
            if slide['notes']:
                section += f"\nDetailed Explanation:\n{slide['notes']}\n"
            
            knowledge_base.append(section)
        
        # Join all sections with clear separators
        return "\n\n" + "="*50 + "\n\n".join(knowledge_base)

    def construct_audio_url(self, filename: str) -> str:
        """
        Construct a full URL for an audio file in the audio folder (test-uploads).
        
        Args:
            filename (str): Name of the audio file
            
        Returns:
            str: Full URL to the audio file
        """
        # Get the endpoint URL without trailing slash
        endpoint = os.getenv("LINODE_ENDPOINT", "").rstrip('/')
        
        # Construct the full URL: endpoint/bucket/folder/filename
        audio_url = f"{endpoint}/{self.bucket_name}/{self.audio_folder}/{filename}"
        
        logger.info(f"Constructed audio URL: {audio_url}")
        return audio_url

    def download_audio_file(self, filename: str, download_path: str = None) -> str:
        """
        Download an audio file from the audio folder (test-uploads) in Linode Object Storage.
        
        Args:
            filename (str): Name of the audio file in the test-uploads folder
            download_path (str): Local path to save the downloaded file
            
        Returns:
            str: Path to the downloaded file
        """
        try:
            logger.info(f"Downloading audio file from {self.audio_folder}: {filename}")
            
            # Construct the key (path in bucket)
            key = f"{self.audio_folder}/{filename}"
            
            # Set default download path if not provided
            if not download_path:
                download_path = f"downloaded_{filename}"
            
            logger.info(f"Downloading from bucket: {self.bucket_name}, key: {key}")
            
            # Check if file exists
            try:
                self.s3_client.head_object(Bucket=self.bucket_name, Key=key)
            except self.s3_client.exceptions.ClientError as e:
                error_code = e.response['Error']['Code']
                if error_code == '404':
                    raise Exception(f"Audio file '{filename}' not found in {self.audio_folder} folder")
                elif error_code == '403':
                    raise Exception(f"Access denied to audio file '{filename}'. Check your credentials.")
                else:
                    raise Exception(f"Error accessing audio file: {str(e)}")
            
            # Download the file
            self.s3_client.download_file(self.bucket_name, key, download_path)
            logger.info(f"Successfully downloaded audio file to: {download_path}")
            
            return download_path
            
        except Exception as e:
            logger.error(f"Error downloading audio file: {str(e)}", exc_info=True)
            raise Exception(f"Failed to download audio file: {str(e)}")

    def list_audio_files(self) -> list:
        """
        List all audio files in the audio folder (test-uploads).
        
        Returns:
            list: List of audio filenames
        """
        try:
            logger.info(f"Listing audio files in {self.audio_folder}")
            
            response = self.s3_client.list_objects_v2(
                Bucket=self.bucket_name,
                Prefix=f"{self.audio_folder}/"
            )
            
            if 'Contents' not in response:
                logger.info(f"No files found in {self.audio_folder}")
                return []
            
            # Extract filenames (remove folder prefix)
            files = [
                obj['Key'].replace(f"{self.audio_folder}/", "")
                for obj in response['Contents']
                if not obj['Key'].endswith('/')  # Exclude folder itself
            ]
            
            logger.info(f"Found {len(files)} audio files")
            return files
            
        except Exception as e:
            logger.error(f"Error listing audio files: {str(e)}", exc_info=True)
            raise Exception(f"Failed to list audio files: {str(e)}")